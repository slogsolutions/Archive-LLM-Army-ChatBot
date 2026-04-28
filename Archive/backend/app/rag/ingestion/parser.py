from __future__ import annotations
import os
import re
from dataclasses import dataclass, field
from typing import List


@dataclass
class ParsedPage:
    page_number: int
    text: str
    heading: str = ""


@dataclass
class ParsedDocument:
    pages: List[ParsedPage]
    file_type: str = ""
    title: str = ""           # document-level title extracted from PDF header

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())

    @property
    def total_pages(self) -> int:
        return len(self.pages)


_KW_STOPWORDS = {
    "the", "and", "for", "with", "this", "that", "from", "have", "been",
    "will", "are", "was", "were", "has", "had", "its", "not", "but",
    "also", "can", "may", "each", "than", "then", "into", "more", "some",
    "used", "use", "using", "which", "when", "where", "their", "they",
    "them", "there", "these", "those", "such", "only", "other", "about",
    "after", "before", "between", "both", "does", "data", "type",
}


def auto_extract_keywords(text: str, top_n: int = 25) -> str:
    """
    Extract top-N keywords from text using simple term frequency.
    Returns a comma-separated string of keywords.
    """
    from collections import Counter
    words = re.findall(r"\b[a-zA-Z]{4,}\b", text.lower())
    words = [w for w in words if w not in _KW_STOPWORDS]
    if not words:
        return ""
    top = [w for w, _ in Counter(words).most_common(top_n)]
    return ", ".join(top)


def extract_metadata(doc, parsed_doc: "ParsedDocument | None" = None) -> dict:
    user_kw = (doc.keywords or "").strip()
    text_for_kw = (doc.corrected_text or doc.ocr_text or "").strip()
    auto_kw = auto_extract_keywords(text_for_kw) if text_for_kw else ""
    combined_kw = ", ".join(filter(None, [user_kw, auto_kw]))

    # Pull title from the parsed document if available
    doc_title = (parsed_doc.title if parsed_doc and parsed_doc.title else "").strip()

    return {
        "branch":          doc.branch_name or "",
        "doc_type":        doc.document_type_name or "",
        "year":            doc.year,
        "section":         doc.section or "",
        "hq_id":           doc.hq_id,
        "unit_id":         doc.unit_id,
        "branch_id":       doc.branch_id,
        "uploaded_by":     doc.uploaded_by,
        "file_name":       doc.file_name,
        "file_type":       doc.file_type or "",
        "min_visible_rank": doc.min_visible_rank if doc.min_visible_rank is not None else 6,
        "keywords":        combined_kw,
        "doc_title":       doc_title,
    }


def parse_document(file_path: str, ocr_text: str = None) -> ParsedDocument:
    """Parse any supported document format into structured pages."""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return _parse_pdf(file_path, ocr_text)
        elif ext in (".docx", ".doc"):
            return _parse_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _parse_xlsx(file_path)
        elif ext == ".csv":
            return _parse_csv(file_path)
        elif ext in (".pptx", ".ppt"):
            return _parse_pptx(file_path)
        elif ext in (".txt", ".text", ".md"):
            return _parse_txt(file_path)
        elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp"):
            return _parse_image(file_path, ocr_text)
        else:
            # Attempt plain text fallback
            return _parse_txt(file_path)
    except Exception as e:
        print(f"[PARSER] Error parsing {ext}: {e}")
        return ParsedDocument(pages=[], file_type=ext.lstrip("."))


_NOTE_RE  = re.compile(r"^note\s*[:\-]", re.IGNORECASE)
_PARA_RE  = re.compile(r"^\d+[\.\)]\s")
_SUB_RE   = re.compile(r"^\([a-z]\)\s", re.IGNORECASE)
_WATER_RE = re.compile(r"^(restricted|confidential|secret|top secret|classified)$", re.IGNORECASE)


def markdown_to_parsed_doc(md: str, method: str = "unknown") -> "ParsedDocument":
    """
    Convert a Markdown string (from Marker/Docling/PyMuPDF) into a
    ParsedDocument where each `##` section becomes one ParsedPage.

    This is the shared entry-point used by pipeline.py after md_parser
    produces the Markdown.  The chunker then receives properly-headed
    pages and can prepend headings to chunk texts without any heuristics.
    """
    from app.rag.ingestion.md_parser import markdown_to_sections

    sections = markdown_to_sections(md)
    if not sections:
        return ParsedDocument(pages=[], file_type="pdf")

    doc_title = sections[0][0] if sections else ""
    pages: List[ParsedPage] = []

    for i, (title, heading, body) in enumerate(sections):
        if not body.strip():
            continue
        pages.append(ParsedPage(
            page_number=i + 1,
            text=body,
            heading=heading,
        ))

    return ParsedDocument(
        pages=[p for p in pages if p.text.strip()],
        file_type="pdf",
        title=doc_title,
    )


def _parse_pdf(file_path: str, ocr_text: str = None) -> ParsedDocument:
    """
    PDF → ParsedDocument via the md_parser strategy cascade:
      Marker → Docling → PyMuPDF (font-aware) → PaddleOCR fallback.

    Each successful strategy produces Markdown; the Markdown is then
    converted into section-per-page ParsedDocument via markdown_to_parsed_doc().
    """
    from app.rag.ingestion.md_parser import convert_to_markdown

    result = convert_to_markdown(file_path=file_path, ocr_text=ocr_text)
    print(
        f"[PARSER] PDF → MD: method={result.method}  "
        f"quality={result.quality:.2f}  chars={len(result.markdown)}"
    )
    if result.warnings:
        for w in result.warnings:
            print(f"[PARSER] Quality warning: {w}")

    if not result.markdown.strip():
        return ParsedDocument(pages=[], file_type="pdf")

    return markdown_to_parsed_doc(result.markdown, method=result.method)


def _parse_docx(file_path: str) -> ParsedDocument:
    from docx import Document as DocxDoc

    doc = DocxDoc(file_path)
    pages: List[ParsedPage] = []
    current_lines: List[str] = []
    current_heading = ""
    page_num = 1

    for para in doc.paragraphs:
        style = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            continue

        if "heading" in style or re.match(r'^(chapter|section|part)\b', text, re.IGNORECASE):
            # Flush current section as a page
            if current_lines:
                pages.append(ParsedPage(
                    page_number=page_num,
                    text="\n".join(current_lines),
                    heading=current_heading,
                ))
                page_num += 1
                current_lines = []
            current_heading = text
            current_lines.append(text)
        else:
            current_lines.append(text)

    # Tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            current_lines.append("\n".join(rows))

    if current_lines:
        pages.append(ParsedPage(
            page_number=page_num,
            text="\n".join(current_lines),
            heading=current_heading,
        ))

    return ParsedDocument(pages=pages, file_type="docx")


def _parse_xlsx(file_path: str) -> ParsedDocument:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    pages: List[ParsedPage] = []

    for idx, sheet in enumerate(wb.worksheets):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            pages.append(ParsedPage(
                page_number=idx + 1,
                text="\n".join(rows),
                heading=sheet.title or f"Sheet {idx + 1}",
            ))

    wb.close()
    return ParsedDocument(pages=pages, file_type="xlsx")


def _parse_csv(file_path: str) -> ParsedDocument:
    import csv

    rows: List[str] = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                if any(c.strip() for c in row):
                    rows.append(" | ".join(row))
    except Exception as e:
        print(f"[PARSER] CSV error: {e}")

    return ParsedDocument(
        pages=[ParsedPage(page_number=1, text="\n".join(rows))],
        file_type="csv",
    )


def _parse_pptx(file_path: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(file_path)
    pages: List[ParsedPage] = []

    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if texts:
            heading = texts[0] if len(texts[0]) < 100 else ""
            pages.append(ParsedPage(
                page_number=i + 1,
                text="\n".join(texts),
                heading=heading,
            ))

    return ParsedDocument(pages=pages, file_type="pptx")


def _parse_txt(file_path: str) -> ParsedDocument:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
    except Exception as e:
        print(f"[PARSER] TXT error: {e}")
        return ParsedDocument(pages=[], file_type="txt")

    if len(text) <= 4000:
        return ParsedDocument(
            pages=[ParsedPage(page_number=1, text=text)],
            file_type="txt",
        )

    # Group paragraphs into logical pages (~2000 chars each)
    paragraphs = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip()]
    pages: List[ParsedPage] = []
    current: List[str] = []
    current_len = 0
    page_num = 1

    for para in paragraphs:
        current.append(para)
        current_len += len(para)
        if current_len >= 2000:
            pages.append(ParsedPage(page_number=page_num, text="\n\n".join(current)))
            page_num += 1
            current = []
            current_len = 0

    if current:
        pages.append(ParsedPage(page_number=page_num, text="\n\n".join(current)))

    return ParsedDocument(pages=pages, file_type="txt")


def _parse_image(file_path: str, ocr_text: str = None) -> ParsedDocument:
    if ocr_text:
        return ParsedDocument(
            pages=[ParsedPage(page_number=1, text=ocr_text)],
            file_type="image",
        )
    return ParsedDocument(pages=[], file_type="image")
